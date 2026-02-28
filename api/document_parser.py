"""
Document parsing utilities for Live Coach context.
Extracts text from PPT, DOCX, PDF, XLSX for LLM context.
"""

import os
import tempfile
from typing import Optional
import logging

logger = logging.getLogger("speechscore.documents")


def extract_text_from_document(file_path: str, filename: str) -> Optional[str]:
    """
    Extract text content from various document formats.
    Returns plain text suitable for LLM context.
    """
    ext = os.path.splitext(filename.lower())[1]
    
    try:
        if ext in ['.pptx', '.ppt']:
            return extract_pptx(file_path)
        elif ext in ['.docx', '.doc']:
            return extract_docx(file_path)
        elif ext == '.pdf':
            return extract_pdf(file_path)
        elif ext in ['.xlsx', '.xls']:
            return extract_xlsx(file_path)
        elif ext == '.txt':
            return extract_txt(file_path)
        else:
            logger.warning(f"Unsupported document type: {ext}")
            return None
    except Exception as e:
        logger.error(f"Error extracting text from {filename}: {e}")
        return None


def extract_pptx(file_path: str) -> str:
    """Extract text from PowerPoint files."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed, cannot parse PPTX")
        return "[PowerPoint file - install python-pptx to extract content]"
    
    prs = Presentation(file_path)
    text_parts = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [f"\n--- Slide {slide_num} ---"]
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
            
            # Handle tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        slide_text.append(row_text)
        
        if len(slide_text) > 1:  # More than just the header
            text_parts.extend(slide_text)
    
    return "\n".join(text_parts)


def extract_docx(file_path: str) -> str:
    """Extract text from Word documents."""
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx not installed, cannot parse DOCX")
        return "[Word document - install python-docx to extract content]"
    
    doc = Document(file_path)
    text_parts = []
    
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())
    
    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                text_parts.append(row_text)
    
    return "\n".join(text_parts)


def extract_pdf(file_path: str) -> str:
    """Extract text from PDF files."""
    try:
        import PyPDF2
    except ImportError:
        logger.warning("PyPDF2 not installed, cannot parse PDF")
        return "[PDF file - install PyPDF2 to extract content]"
    
    text_parts = []
    
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                text_parts.append(f"\n--- Page {page_num} ---")
                text_parts.append(text.strip())
    
    return "\n".join(text_parts)


def extract_xlsx(file_path: str) -> str:
    """Extract text from Excel spreadsheets."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.warning("openpyxl not installed, cannot parse XLSX")
        return "[Excel spreadsheet - install openpyxl to extract content]"
    
    wb = load_workbook(file_path, read_only=True, data_only=True)
    text_parts = []
    
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        sheet_text = [f"\n--- Sheet: {sheet_name} ---"]
        
        for row in sheet.iter_rows(max_row=100):  # Limit rows
            row_values = [str(cell.value).strip() for cell in row if cell.value is not None]
            if row_values:
                sheet_text.append(" | ".join(row_values))
        
        if len(sheet_text) > 1:
            text_parts.extend(sheet_text)
    
    wb.close()
    return "\n".join(text_parts)


def extract_txt(file_path: str) -> str:
    """Extract text from plain text files."""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()


def summarize_for_context(text: str, max_chars: int = 8000) -> str:
    """
    Prepare document text for LLM context.
    Truncates if too long, keeping beginning and end.
    """
    if not text:
        return ""
    
    if len(text) <= max_chars:
        return text
    
    # Keep first 60% and last 40%
    first_part = int(max_chars * 0.6)
    last_part = max_chars - first_part - 50  # Leave room for truncation notice
    
    return (
        text[:first_part] + 
        "\n\n[... content truncated for length ...]\n\n" + 
        text[-last_part:]
    )
